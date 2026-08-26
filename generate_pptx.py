"""BinlogTrace 功能介绍 PPT 生成器."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, zipfile, re

PALETTE = {
    "dominant": RGBColor(0x1B, 0x2A, 0x4A),
    "support1": RGBColor(0x7A, 0xC7, 0xE3),
    "support2": RGBColor(0xF3, 0xF5, 0xF8),
    "accent":   RGBColor(0xF0, 0x96, 0x2E),
    "ink":      RGBColor(0x21, 0x21, 0x21),
    "muted":    RGBColor(0x6B, 0x6B, 0x6B),
    "green":    RGBColor(0x28, 0xA7, 0x45),
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
}
FONT_HEAD = "Calibri"; FONT_BODY = "Calibri"
FONT_HEAD_EA = "Microsoft YaHei"; FONT_BODY_EA = "Microsoft YaHei"

CHAR = {"title": 60, "sub": 100, "bullet": 95, "foot": 80, "label": 30}

def guard(t, key):
    lim = CHAR[key]
    if len(t) <= lim: return t
    cut = t[:lim-1]; sp = cut.rfind(" ")
    return (cut[:sp] if sp > 0 else cut) + "…"

def add_text(slide, text, x, y, w, h, *, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, font=None, font_ea=None, budget_key=None):
    if budget_key: text = guard(text, budget_key)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font or FONT_BODY
    r.font.size = Pt(size); r.font.bold = bold
    if color is not None: r.font.color.rgb = color
    from pptx.oxml.ns import qn
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {"typeface": font_ea or FONT_BODY_EA}); rPr.append(ea)
    else: ea.set("typeface", font_ea or FONT_BODY_EA)
    return tb

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def motif_bar(slide): add_rect(slide, 0, 0, 0.18, 7.5, PALETTE["accent"])

def title_bar(slide, title):
    slide2 = slide
    motif_bar(slide2)
    add_text(slide2, title, 0.6, 0.4, 12.1, 0.9, size=30, bold=True,
             color=PALETTE["dominant"], font=FONT_HEAD, font_ea=FONT_HEAD_EA, budget_key="title")
    add_rect(slide2, 0.6, 1.3, 12.1, 0.03, PALETTE["accent"])

def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        try:
            slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
            return
        except Exception as e:
            add_text(slide, "[图片加载失败]", x, y, w, 0.5, color=PALETTE["muted"], size=12)
            return
    add_text(slide, "[缺少截图: "+os.path.basename(path)+"]", x, y, w, 0.5, color=PALETTE["muted"], size=12)

def save_pptx(prs, path):
    prs.save(path)
    tmp = path + ".tmp"
    try:
        with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w") as zout:
            for item in zin.infolist():
                if item.filename.lower().startswith("docprops/thumbnail"): continue
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    c = data.decode("utf-8")
                    c = re.sub(r'<Override[^>]*PartName="/docProps/thumbnail[^"]*"[^>]*/?>', '', c, flags=re.I)
                    data = c.encode("utf-8")
                if item.filename == "_rels/.rels":
                    c = data.decode("utf-8")
                    c = re.sub(r'<Relationship[^>]*Target="docProps/thumbnail[^"]*"[^>]*/?>', '', c, flags=re.I)
                    data = c.encode("utf-8")
                zout.writestr(item, data)
        os.replace(tmp, path)
    except Exception as e:
        print("warning: strip thumbnail", e)
        if os.path.exists(tmp):
            try: os.remove(tmp)
            except: pass

SLIDES_DIR = "outputs"

# ---------------- Slide builders ----------------
def s_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0,0,13.333,7.5, PALETTE["dominant"])
    add_rect(s, 0,5.5,13.333,0.06, PALETTE["accent"])
    add_text(s, "BinlogTrace", 0.9,2.3,11.5,1.5, size=50,bold=True,color=PALETTE["white"],
             font=FONT_HEAD,font_ea=FONT_HEAD_EA,align=PP_ALIGN.CENTER,budget_key="title")
    add_text(s, "MySQL 数据追踪与回滚脚本生成工具", 0.9,3.9,11.5,0.8, size=22,
             color=PALETTE["support1"],font=FONT_HEAD,font_ea=FONT_HEAD_EA,align=PP_ALIGN.CENTER,budget_key="sub")
    add_text(s, "基于 binlog 实时追踪变更 / 一键生成回滚 SQL / 双引擎架构（Workerman + TypePHP）",
             0.9,5.8,11.5,0.6, size=12, color=PALETTE["support1"], align=PP_ALIGN.CENTER, budget_key="sub")
    return s

def s_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "核心功能")
    items = [
        ("连接配置", "真实 MySQL / 演示模式切换，测试连接 + 权限预检（SELECT / REPLICATION）"),
        ("追踪工单", "选库/表 + 时间窗口 + INSERT/UPDATE/DELETE 类型筛选"),
        ("变更列表", "类型、库.表、主键、变更列一目了然，分页浏览 100 条/页"),
        ("变更明细·差异", "列名 / 旧值 / 新值 字段级对比，清晰定位变化"),
        ("回滚脚本", "单行「生成该行回滚」+ 多选「生成回滚脚本」，支持独立/整体事务"),
        ("双引擎架构", "agent-workerman（纯 PHP）与 agent（TypePHP 编译版）可对比"),
    ]
    y = 1.7
    for i, (t, d) in enumerate(items, 1):
        add_rect(s, 0.7, y, 0.5, 0.5, PALETTE["dominant"])
        add_text(s, str(i), 0.7, y, 0.5, 0.5, size=14, bold=True, color=PALETTE["white"],
                 align=PP_ALIGN.CENTER, budget_key="label")
        add_text(s, t, 1.4, y, 3.0, 0.5, size=16, bold=True, color=PALETTE["dominant"], budget_key="label")
        add_text(s, d, 4.6, y, 8.3, 0.5, size=13, color=PALETTE["ink"], budget_key="bullet")
        y += 0.86
    return s

def s_flow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "使用流程")
    add_img(s, os.path.join(SLIDES_DIR, "01-connect-page.png"), 0.6, 1.6, 5.6, 5.4)
    add_img(s, os.path.join(SLIDES_DIR, "02-trace-config.png"), 6.6, 1.6, 6.1, 5.4)
    add_text(s, "① 连接配置页：选库/类型 + 测试连接", 0.6, 6.95, 5.6, 0.4, size=11, color=PALETTE["muted"], budget_key="label")
    add_text(s, "② 追踪工单：选库 + 时间窗口 + 类型", 6.6, 6.95, 6.1, 0.4, size=11, color=PALETTE["muted"], budget_key="label")
    return s

def s_result(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "变更列表 · 真实数据展示")
    add_img(s, os.path.join(SLIDES_DIR, "03-result-page.png"), 0.7, 1.5, 12, 5.5)
    add_text(s, "库.表显示真实表名（shop.orders / shop.payments / shop.users）· 主键有值（id=1000）· 已拉取大量变更",
             0.7, 6.9, 12.0, 0.4, size=12, color=PALETTE["muted"], budget_key="sub")
    return s

def s_detail(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "变更明细 · 字段级差异")
    add_img(s, os.path.join(SLIDES_DIR, "04-change-detail.png"), 0.7, 1.5, 6.8, 5.5)
    add_text(s, "列名 / 旧值 / 新值 对比表格", 0.7, 6.9, 6.8, 0.4, size=12, color=PALETTE["muted"])
    bullets = [
        "主键定位：id=1001",
        "UPDATE 展示旧/新值：amount 69.47 → 169.47",
        "支持「生成该行回滚」",
        "事务号、binlog 位置标注",
    ]
    y = 1.8
    for b in bullets:
        add_rect(s, 8.0, y, 0.4, 0.4, PALETTE["accent"])
        add_text(s, "•", 8.0, y, 0.4, 0.4, size=14, bold=True, color=PALETTE["white"], align=PP_ALIGN.CENTER, budget_key="label")
        add_text(s, b, 8.6, y, 4.3, 0.7, size=14, color=PALETTE["ink"], budget_key="bullet")
        y += 0.85
    return s

def s_rollback(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "回滚脚本 · 单行 / 独立事务")
    add_img(s, os.path.join(SLIDES_DIR, "05-row-rollback.png"), 0.6, 1.5, 6.4, 5.5)
    add_img(s, os.path.join(SLIDES_DIR, "08-independent-rollback.png"), 7.3, 1.5, 5.5, 5.5)
    add_text(s, "单行回滚 SQL 预览（含注释、源操作、WHERE 定位）", 0.6, 6.9, 6.4, 0.4, size=11, color=PALETTE["muted"])
    add_text(s, "独立回滚：每条变更单独事务（3 事务 · 3 语句）", 7.3, 6.9, 5.5, 0.4, size=11, color=PALETTE["muted"])
    return s

def s_multirollback(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "批量回滚 · 多选 + 独立回滚切换")
    add_img(s, os.path.join(SLIDES_DIR, "06-multi-select.png"), 0.6, 1.5, 6.4, 5.4)
    add_img(s, os.path.join(SLIDES_DIR, "07-multi-rollback.png"), 7.3, 1.5, 5.4, 5.4)
    add_text(s, "勾选多条 → 生成回滚脚本 (3)", 0.6, 6.9, 6.4, 0.4, size=11, color=PALETTE["muted"])
    add_text(s, "整体事务：1 事务 · 3 语句 · 34 行", 7.3, 6.9, 5.4, 0.4, size=11, color=PALETTE["muted"])
    return s

def s_arch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "双引擎架构对比")
    cols = [
        ("agent-workerman（纯 PHP）", ["Workerman 5 事件驱动", "krowinski 子进程解析 binlog", "HTTP 模式（fetch POST）", "PDO 适配器（修复认证）"]),
        ("agent（TypePHP 编译版）", ["原生二进制，直接 TCP 流解析", "AsyncClient 手写协议", "无子进程 / 无文件轮询", "编译生成 binlog_agent"]),
    ]
    x = 0.7
    for title, bullets in cols:
        add_rect(s, x, 1.7, 5.8, 0.7, PALETTE["dominant"])
        add_text(s, title, x, 1.78, 5.8, 0.5, size=15, bold=True, color=PALETTE["white"], align=PP_ALIGN.CENTER, budget_key="sub")
        y = 2.6
        for b in bullets:
            add_text(s, "• " + b, x+0.3, y, 5.2, 0.7, size=13, color=PALETTE["ink"], budget_key="bullet")
            y += 0.75
        x += 6.3
    add_text(s, "两版可对比速度与解析路径；演示模式数据固定，账号密码任意填即可联调。",
             0.7, 6.4, 12.0, 0.6, size=13, color=PALETTE["muted"], budget_key="sub")
    return s

def s_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0,0,13.333,7.5, PALETTE["dominant"])
    add_text(s, "谢谢观看", 1.0, 3.0, 11.3, 1.5, size=40, bold=True, color=PALETTE["white"],
             align=PP_ALIGN.CENTER, font=FONT_HEAD, font_ea=FONT_HEAD_EA, budget_key="title")
    add_text(s, "BinlogTrace — 让 MySQL 每个变更都有迹可循、可回滚", 1.0, 4.8, 11.3, 0.6,
             size=16, color=PALETTE["support1"], align=PP_ALIGN.CENTER, budget_key="sub")
    return s

def main():
    prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    s_cover(prs)
    s_overview(prs)
    s_flow(prs)
    s_result(prs)
    s_detail(prs)
    s_rollback(prs)
    s_multirollback(prs)
    s_arch(prs)
    s_closing(prs)
    import os as mm; mm.makedirs("outputs", exist_ok=True)
    out = "outputs/BinlogTrace-功能介绍.pptx"
    save_pptx(prs, out)
    print("wrote", out)

if __name__ == "__main__":
    main()